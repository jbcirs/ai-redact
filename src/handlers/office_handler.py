#!/usr/bin/env python3
"""
office_handler.py — Kind A convert handler for Word (.docx) and
PowerPoint (.pptx).

Tier-1 pure-Python conversion (docs/plans/format-support-plan.md §3.4):
content is read with python-docx / python-pptx and rendered into a clean,
simplified PDF with PyMuPDF. The router then runs the proven PDF
redaction/verification pipeline over the result. Layout is simplified,
NOT pixel-identical — the goal is content-faithful redaction prep, not
document reproduction.

Safety rule: anything this converter knows exists but cannot carry over
(text boxes, SmartArt/AlternateContent, charts, unreadable images) is
OMITTED from the output — the safe direction, omitted content cannot
leak — but is always COUNTED in info["dropped_elements"] so loss is
never silent.

Local only. No network. No subprocesses.
"""

from __future__ import annotations

import sys
import zipfile
from pathlib import Path

try:
    import fitz  # PyMuPDF
except ImportError:  # pragma: no cover
    sys.exit(
        "PyMuPDF is not installed.\n"
        "Run:  pip install -r requirements.txt"
    )

try:  # imported as part of the handlers package (normal router path)
    from handlers.common import UnsupportedFormatError
    from handlers.pdf_render import PdfFlow as _PdfFlow
except ImportError:  # executed directly: python src/handlers/office_handler.py
    from common import UnsupportedFormatError
    from pdf_render import PdfFlow as _PdfFlow

SUPPORTED_EXTENSIONS = {".docx", ".pptx"}


# ---------------------------------------------------------------------------
# .docx conversion
# ---------------------------------------------------------------------------


def _docx_count_unconvertible(input_path: Path) -> tuple[int, list[str]]:
    """Count content python-docx cannot see, straight off the raw XML.

    BlockItemContainer walking never reaches text boxes (w:txbxContent) or
    SmartArt/fallback pairs (mc:AlternateContent), so their presence is
    detected by counting opening tags in document.xml and every header/
    footer part. May overcount (Choice/Fallback duplicate their content) —
    overcounting dropped elements is the safe direction."""
    dropped = 0
    notes: list[str] = []
    textboxes = 0
    alternates = 0
    ole_objects = 0
    try:
        with zipfile.ZipFile(input_path) as zf:
            names = [
                n
                for n in zf.namelist()
                if n == "word/document.xml"
                or (n.startswith("word/header") and n.endswith(".xml"))
                or (n.startswith("word/footer") and n.endswith(".xml"))
            ]
            for name in names:
                xml = zf.read(name).decode("utf-8", "replace")
                textboxes += xml.count("<w:txbxContent")
                alternates += xml.count("<mc:AlternateContent")
                ole_objects += xml.count("<w:object")
    except Exception:
        notes.append("could not inspect raw XML for text boxes/SmartArt")
        return 0, notes
    if textboxes:
        dropped += textboxes
        notes.append(
            f"{textboxes} text box(es) not converted (counted as dropped)"
        )
    if alternates:
        dropped += alternates
        notes.append(
            f"{alternates} AlternateContent element(s) (SmartArt/shapes) "
            "not converted (counted as dropped)"
        )
    if ole_objects:
        dropped += ole_objects
        notes.append(
            f"{ole_objects} embedded OLE object(s) not converted "
            "(counted as dropped)"
        )
    return dropped, notes


def _docx_footnotes(input_path: Path) -> tuple[list[str], int]:
    """Extract footnote/endnote text straight from the zip (python-docx has
    no footnote API). Returns (texts, dropped) — dropped counts note parts
    that exist but could not be parsed. Never raises."""
    texts: list[str] = []
    dropped = 0
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    for part_name, tag in (
        ("word/footnotes.xml", "footnote"),
        ("word/endnotes.xml", "endnote"),
    ):
        try:
            with zipfile.ZipFile(input_path) as zf:
                if part_name not in zf.namelist():
                    continue
                blob = zf.read(part_name)
            from lxml import etree

            root = etree.fromstring(blob)
            for note in root.iter(f"{{{w_ns}}}{tag}"):
                # Skip separator/continuation pseudo-notes.
                if note.get(f"{{{w_ns}}}type"):
                    continue
                text = "".join(
                    t.text or "" for t in note.iter(f"{{{w_ns}}}t")
                ).strip()
                if text:
                    texts.append(text)
        except Exception:
            dropped += 1  # the part exists but is unreadable — count it
    return texts, dropped


def _docx_to_pdf(input_path: Path, options: dict) -> tuple[bytes, dict]:
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise UnsupportedFormatError(
            "python-docx is not installed — run "
            "'pip install -r requirements.txt' to convert .docx files."
        ) from exc

    try:
        doc = Document(str(input_path))
    except Exception as exc:
        raise UnsupportedFormatError(
            f"could not open {input_path.name} as a .docx file "
            f"({exc}). If it is a legacy .doc, save it as .docx first."
        ) from exc

    flow = _PdfFlow()
    dropped, notes = _docx_count_unconvertible(input_path)

    def emit_paragraph(para: Paragraph, size: float = 11.0) -> None:
        text = para.text
        if text and text.strip():
            flow.add_text(text, size=size)

    def emit_table(table: Table) -> None:
        for row in table.rows:
            cells = [c.text.replace("\n", " / ").strip() for c in row.cells]
            flow.add_text(" | ".join(cells), size=10.0)
        flow.add_gap()

    def emit_container(container) -> None:
        """Emit a header/footer's paragraphs and tables in order."""
        for para in container.paragraphs:
            emit_paragraph(para, size=10.0)
        for table in container.tables:
            emit_table(table)

    # Headers — all three per-section variants, skipping linked duplicates.
    header_attrs = ("header", "first_page_header", "even_page_header")
    footer_attrs = ("footer", "first_page_footer", "even_page_footer")
    for section in doc.sections:
        for attr in header_attrs:
            hf = getattr(section, attr, None)
            if hf is not None and not hf.is_linked_to_previous:
                emit_container(hf)
    flow.add_gap(10)

    # Body: walk document.element.body children so paragraphs and tables
    # come out in true document order.
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            emit_paragraph(Paragraph(child, doc))
        elif child.tag == qn("w:tbl"):
            emit_table(Table(child, doc))
        # w:sectPr and friends carry no user text — nothing to emit.

    # Embedded images — simplified placement after the paragraph flow.
    image_blobs: list[bytes] = []
    for part in doc.part.related_parts.values():
        content_type = getattr(part, "content_type", "") or ""
        if content_type.startswith("image/"):
            try:
                image_blobs.append(part.blob)
            except Exception:
                dropped += 1
    placed = 0
    for blob in image_blobs:
        flow.add_gap(8)
        if flow.add_image(blob):
            placed += 1
        else:
            dropped += 1
    if placed:
        notes.append(f"{placed} embedded image(s) placed after text flow")
    if placed < len(image_blobs):
        notes.append(
            f"{len(image_blobs) - placed} image(s) could not be rendered "
            "(counted as dropped)"
        )

    # Footnotes/endnotes — read straight from the zip.
    footnotes, fn_dropped = _docx_footnotes(input_path)
    dropped += fn_dropped
    if footnotes:
        flow.add_gap(10)
        flow.add_text("Footnotes:", size=10.0)
        for i, text in enumerate(footnotes, 1):
            flow.add_text(f"{i}. {text}", size=10.0)
        notes.append("footnotes included")
    if fn_dropped:
        notes.append(
            f"{fn_dropped} footnote/endnote part(s) unreadable "
            "(counted as dropped)"
        )

    # Footers last.
    flow.add_gap(10)
    for section in doc.sections:
        for attr in footer_attrs:
            hf = getattr(section, attr, None)
            if hf is not None and not hf.is_linked_to_previous:
                emit_container(hf)

    if flow.substituted_chars:
        notes.append(
            f"{flow.substituted_chars} character(s) outside the plain "
            "font's repertoire replaced with '?'"
        )
    notes.insert(0, "layout simplified; content-faithful")
    info = {
        "converter": "tier1-python-docx",
        "dropped_elements": dropped,
        "notes": notes,
    }
    return flow.to_bytes(), info


# ---------------------------------------------------------------------------
# .pptx conversion
# ---------------------------------------------------------------------------


def _pptx_to_pdf(input_path: Path, options: dict) -> tuple[bytes, dict]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise UnsupportedFormatError(
            "python-pptx is not installed — run "
            "'pip install -r requirements.txt' to convert .pptx files."
        ) from exc

    try:
        prs = Presentation(str(input_path))
    except Exception as exc:
        raise UnsupportedFormatError(
            f"could not open {input_path.name} as a .pptx file "
            f"({exc}). If it is a legacy .ppt, save it as .pptx first."
        ) from exc

    flow = _PdfFlow()
    dropped = 0
    notes: list[str] = []
    notes_included = False
    a_t = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"

    def shape_has_hidden_text(shape) -> bool:
        """True when the shape's XML holds text we did not extract
        (SmartArt and other graphicFrames keep text in a:t runs)."""
        try:
            return any(
                (t.text or "").strip()
                for t in shape._element.iter(a_t)
            )
        except Exception:
            return True  # unreadable — assume the worst, count it

    def emit_shape(shape) -> None:
        nonlocal dropped
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                emit_shape(sub)
            return
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [
                    c.text.replace("\n", " / ").strip() for c in row.cells
                ]
                flow.add_text(" | ".join(cells), size=10.0)
            flow.add_gap()
            return
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text
            if text and text.strip():
                flow.add_text(text)
                flow.add_gap(4)
            return
        if getattr(shape, "has_chart", False):
            # Chart text (title/labels) lives in a separate part we do not
            # render — count the chart as dropped, never silently.
            dropped += 1
            return
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(
            shape, "image"
        ):
            try:
                blob = shape.image.blob
            except Exception:
                dropped += 1
                return
            if not flow.add_image(blob):
                dropped += 1
            return
        if shape_has_hidden_text(shape):
            dropped += 1

    slide_count = 0
    for slide in prs.slides:
        if slide_count:
            flow.new_page()  # one PDF page per slide
        slide_count += 1
        for shape in slide.shapes:
            emit_shape(shape)
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text and notes_text.strip():
                flow.add_gap(10)
                flow.add_text("Notes:", size=10.0)
                flow.add_text(notes_text, size=10.0)
                notes_included = True

    notes.append("layout simplified; content-faithful")
    notes.append(f"{slide_count} slides")
    if notes_included:
        notes.append("speaker notes included")
    if flow.substituted_chars:
        notes.append(
            f"{flow.substituted_chars} character(s) outside the plain "
            "font's repertoire replaced with '?'"
        )
    info = {
        "converter": "tier1-python-pptx",
        "dropped_elements": dropped,
        "notes": notes,
    }
    return flow.to_bytes(), info


# ---------------------------------------------------------------------------
# Handler entry point (Kind A contract)
# ---------------------------------------------------------------------------


def to_pdf(input_path: Path, options: dict) -> tuple[bytes, dict]:
    """Convert a .docx or .pptx to simplified-PDF bytes (Tier 1).

    Returns (pdf_bytes, info) per the Kind A handler contract:
    info = {"converter", "dropped_elements", "notes"}.
    """
    input_path = Path(input_path)
    ext = input_path.suffix.lower()
    if ext == ".docx":
        return _docx_to_pdf(input_path, options or {})
    if ext == ".pptx":
        return _pptx_to_pdf(input_path, options or {})
    raise UnsupportedFormatError(
        f"office_handler cannot convert '{ext}' files — supported: "
        ".docx, .pptx. Legacy .doc/.ppt must be re-saved in the modern "
        "format first."
    )


# ---------------------------------------------------------------------------
# Standalone smoke test
# ---------------------------------------------------------------------------


def _smoke_test() -> int:
    import tempfile

    sys.path.insert(
        0, str(Path(__file__).resolve().parents[2] / "tests")
    )
    from make_office_fixtures import (
        PLANTED_EMAIL,
        PLANTED_PHONE,
        PLANTED_SSN,
        SURVIVOR,
        make_docx,
        make_pptx,
    )

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)

        # --- .docx: conversion must carry body, table, and header text ---
        docx_path = make_docx(tmp_dir)
        pdf_bytes, info = to_pdf(docx_path, {})
        pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = "\n".join(page.get_text() for page in pdf)
        pdf.close()
        assert PLANTED_EMAIL in text, "docx: planted email missing from PDF"
        assert "Header contact:" in text, "docx: header text missing"
        assert PLANTED_SSN in text, "docx: table SSN missing"
        assert SURVIVOR in text, "docx: dollar amount missing"
        assert info["converter"] == "tier1-python-docx"
        assert isinstance(info["dropped_elements"], int)
        print(f"docx info: {info}")

        # --- .pptx: slides, table, and speaker notes must carry over ---
        pptx_path = make_pptx(tmp_dir)
        pdf_bytes, info = to_pdf(pptx_path, {})
        pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
        text = "\n".join(page.get_text() for page in pdf)
        page_count = pdf.page_count
        pdf.close()
        assert PLANTED_EMAIL in text, "pptx: planted email missing"
        assert PLANTED_PHONE in text, "pptx: speaker-note phone missing"
        assert SURVIVOR in text, "pptx: table dollar amount missing"
        assert page_count >= 2, "pptx: expected one page per slide"
        assert info["converter"] == "tier1-python-pptx"
        print(f"pptx info: {info}")

    print("office_handler smoke test: PASS")
    return 0


if __name__ == "__main__":
    try:
        sys.exit(_smoke_test())
    except AssertionError as exc:
        print(f"office_handler smoke test: FAIL — {exc}", file=sys.stderr)
        sys.exit(1)
