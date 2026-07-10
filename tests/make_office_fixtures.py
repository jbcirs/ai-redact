#!/usr/bin/env python3
"""
make_office_fixtures.py — generate fake-PII Office fixtures for handler tests.

Writes into the directory given as argv[1]:
  fixture.docx  — paragraphs + table + header with planted identifiers
  fixture.pptx  — 2 slides (title/body, table) + speaker notes
  fixture.xlsx  — 2 sheets (one hidden), comment, defined name, core props

Planted values (per docs/plans/handler-spec.md):
  email  planted.email@example.com
  phone  (555) 010-9999
  ssn    000-55-4444
  name   Casey Plantedname
  MUST-SURVIVE financial string: $12,345.67

Pure Python, no network; deps from requirements.txt only.
"""

from __future__ import annotations

import sys
from pathlib import Path

PLANTED_EMAIL = "planted.email@example.com"
PLANTED_PHONE = "(555) 010-9999"
PLANTED_SSN = "000-55-4444"
PLANTED_NAME = "Casey Plantedname"
SURVIVOR = "$12,345.67"


def make_docx(out_dir: Path) -> Path:
    """Word fixture: body paragraphs, a table holding the SSN, a header
    holding the email, and the must-survive dollar amount."""
    from docx import Document

    path = Path(out_dir) / "fixture.docx"
    doc = Document()

    # Header (section 0) — carries the planted email with a distinctive
    # marker so tests can prove header text made it into the conversion.
    header = doc.sections[0].header
    header.paragraphs[0].text = f"Header contact: {PLANTED_EMAIL}"

    doc.add_heading("Quarterly Account Review", level=1)
    doc.add_paragraph(
        f"Prepared by {PLANTED_NAME}. Direct line {PLANTED_PHONE}."
    )
    doc.add_paragraph(f"Email inquiries to {PLANTED_EMAIL} within 5 days.")
    doc.add_paragraph(f"Closing balance for the period: {SURVIVOR} (verified).")

    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Field"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Taxpayer SSN"
    table.cell(1, 1).text = PLANTED_SSN

    footer = doc.sections[0].footer
    footer.paragraphs[0].text = "Confidential — internal use only"

    doc.save(path)
    print(f"wrote {path}")
    return path


def make_pptx(out_dir: Path) -> Path:
    """PowerPoint fixture: slide 1 title/body with planted values and
    speaker notes holding the phone; slide 2 a table with the survivor."""
    from pptx import Presentation
    from pptx.util import Inches

    path = Path(out_dir) / "fixture.pptx"
    prs = Presentation()

    # Slide 1 — title + content layout.
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = f"Account Review — {PLANTED_NAME}"
    body = slide1.placeholders[1].text_frame
    body.text = f"Contact: {PLANTED_EMAIL}"
    body.add_paragraph().text = f"SSN on file: {PLANTED_SSN}"
    notes = slide1.notes_slide.notes_text_frame
    notes.text = f"Presenter: call back at {PLANTED_PHONE} before Friday."

    # Slide 2 — blank layout with a table.
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    table = slide2.shapes.add_table(
        2, 2, Inches(1), Inches(1), Inches(8), Inches(2)
    ).table
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = PLANTED_NAME
    table.cell(1, 0).text = "Balance"
    table.cell(1, 1).text = SURVIVOR

    prs.save(path)
    print(f"wrote {path}")
    return path


def make_xlsx(out_dir: Path) -> Path:
    """Excel fixture: visible sheet + hidden sheet, planted values across
    cells (one posing as a formula result — a plain value is what
    data_only=True would read anyway), a cell comment with the planted
    name, a defined-name constant with the email, PII in core properties,
    and the must-survive dollar amount."""
    from openpyxl import Workbook
    from openpyxl.comments import Comment
    from openpyxl.workbook.defined_name import DefinedName

    path = Path(out_dir) / "fixture.xlsx"
    wb = Workbook()

    ws1 = wb.active
    ws1.title = "Accounts"
    ws1["A1"] = "Field"
    ws1["B1"] = "Value"
    ws1["A2"] = "Owner email"
    ws1["B2"] = PLANTED_EMAIL
    ws1["A3"] = "Phone"
    ws1["B3"] = PLANTED_PHONE
    ws1["A4"] = "Balance"
    ws1["B4"] = SURVIVOR  # must survive redaction untouched
    # "Formula result": with data_only=True only the cached value is read,
    # so a plain value stands in for a computed one.
    ws1["A5"] = "Derived contact (formula result)"
    ws1["B5"] = f"reach: {PLANTED_EMAIL}"
    ws1["A6"].comment = Comment(
        f"Reviewed by {PLANTED_NAME} on 2026-07-01", "Reviewer"
    )

    ws2 = wb.create_sheet("Archive")
    ws2.sheet_state = "hidden"
    ws2["A1"] = "Legacy record"
    ws2["A2"] = PLANTED_SSN
    ws2["B2"] = f"backup mail {PLANTED_EMAIL}"

    # Defined-name string constant carrying PII.
    wb.defined_names["ContactEmail"] = DefinedName(
        "ContactEmail", attr_text=f'"{PLANTED_EMAIL}"'
    )

    # Core properties carrying PII (must be cleared by the handler).
    wb.properties.creator = PLANTED_EMAIL
    wb.properties.title = f"Workbook of {PLANTED_NAME}"

    wb.save(path)
    print(f"wrote {path}")
    return path


def main(argv: list[str]) -> int:
    if len(argv) < 2:
        print("usage: make_office_fixtures.py <output-dir>", file=sys.stderr)
        return 2
    out_dir = Path(argv[1])
    out_dir.mkdir(parents=True, exist_ok=True)
    make_docx(out_dir)
    make_pptx(out_dir)
    make_xlsx(out_dir)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
